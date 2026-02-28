"""
Multi-Format Export Engine
Generates downloadable files (PDF / DOCX / PPTX / XLSX) from a SynthesisReport
and optional FinancialTimeline data.
"""

from __future__ import annotations

import io
import re
import textwrap
from datetime import datetime
from typing import Any, Optional

from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

router = APIRouter()

# ─── Request schema ──────────────────────────────────────────────────────────

class ExportRequest(BaseModel):
    format: str                          # "pdf" | "docx" | "pptx" | "xlsx"
    session_id: str
    report: dict[str, Any]               # Full SynthesisReport JSON
    timeline: Optional[list[dict]] = None   # ComputedMonthPoint[] – for xlsx

# ─── Helpers ─────────────────────────────────────────────────────────────────

_CITE_RE = re.compile(r"\[\d+\]")

def strip_citations(text: str) -> str:
    return _CITE_RE.sub("", text or "").strip()

def safe(report: dict, key: str, default: str = "") -> str:
    return strip_citations(str(report.get(key) or default))

def references(report: dict) -> list[dict]:
    return report.get("references") or []

def findings(report: dict) -> list[dict]:
    return report.get("key_findings") or []

def swot_items(report: dict, key: str) -> list[str]:
    return [strip_citations(i) for i in (report.get("swot") or {}).get(key, [])]

def business_name(report: dict) -> str:
    name = report.get("business_name") or report.get("session_id") or "Your Business"
    return strip_citations(str(name))

BRAND = "GuggleSimulation"
BRAND_COLOUR = (99, 102, 241)  # indigo-500

# ─── PDF ─────────────────────────────────────────────────────────────────────

def _build_pdf(report: dict) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        HRFlowable,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=22 * mm,
        bottomMargin=18 * mm,
    )

    styles = getSampleStyleSheet()
    indigo = colors.Color(*[x / 255 for x in BRAND_COLOUR])
    dark = colors.HexColor("#1a1a2e")
    light_gray = colors.HexColor("#6b7280")
    amber = colors.HexColor("#f59e0b")

    H1 = ParagraphStyle("H1", parent=styles["Normal"],
                         fontSize=22, fontName="Helvetica-Bold",
                         textColor=indigo, spaceAfter=6)
    H2 = ParagraphStyle("H2", parent=styles["Normal"],
                         fontSize=13, fontName="Helvetica-Bold",
                         textColor=colors.HexColor("#1e293b"), spaceBefore=14, spaceAfter=4)
    BODY = ParagraphStyle("Body", parent=styles["Normal"],
                           fontSize=10, leading=15, textColor=colors.HexColor("#374151"))
    META = ParagraphStyle("Meta", parent=styles["Normal"],
                           fontSize=8.5, textColor=light_gray, spaceAfter=6)
    PITCH = ParagraphStyle("Pitch", parent=styles["Normal"],
                            fontSize=11, fontName="Helvetica-Oblique", leading=16,
                            textColor=colors.HexColor("#312e81"),
                            backColor=colors.HexColor("#eef2ff"),
                            leftIndent=10, rightIndent=10, spaceAfter=6)
    CITE = ParagraphStyle("Cite", parent=styles["Normal"],
                           fontSize=8, leading=12, textColor=light_gray)

    name = business_name(report)
    refs = references(report)
    story = []

    # ── Cover ──────────────────────────────────────────────────────────────
    story.append(Spacer(1, 20 * mm))
    story.append(Paragraph(BRAND, META))
    story.append(Paragraph(f"Executive Blueprint — {name}", H1))
    story.append(Paragraph(
        f"Generated {datetime.now().strftime('%B %d, %Y')}  ·  "
        f"Confidence score: <b>{report.get('confidence_score', 0):.0f}%</b>",
        META,
    ))
    story.append(HRFlowable(width="100%", color=indigo, thickness=1, spaceAfter=10))

    # ── Investor Pitch ─────────────────────────────────────────────────────
    if report.get("investor_pitch"):
        story.append(Paragraph("Investor Pitch", H2))
        story.append(Paragraph(f"&ldquo;{safe(report, 'investor_pitch')}&rdquo;", PITCH))

    # ── Overall Summary ────────────────────────────────────────────────────
    if report.get("overall_summary"):
        story.append(Paragraph("Overview", H2))
        story.append(Paragraph(safe(report, "overall_summary"), BODY))

    # ── Executive Summary ──────────────────────────────────────────────────
    story.append(Paragraph("Executive Summary", H2))
    story.append(Paragraph(safe(report, "executive_summary"), BODY))

    # ── Key Findings ───────────────────────────────────────────────────────
    kf = findings(report)
    if kf:
        story.append(Paragraph("Key Findings", H2))
        for f in kf:
            bullet = f"• [{f.get('confidence_level','–')}] {strip_citations(f.get('claim', ''))}"
            story.append(Paragraph(bullet, BODY))

    # ── SWOT ───────────────────────────────────────────────────────────────
    story.append(Paragraph("SWOT Analysis", H2))
    swot_data = [
        ["✅ Strengths", "⚠️ Weaknesses"],
        [
            "\n".join(f"• {x}" for x in swot_items(report, "strengths")),
            "\n".join(f"• {x}" for x in swot_items(report, "weaknesses")),
        ],
        ["🌍 Opportunities", "🔴 Threats"],
        [
            "\n".join(f"• {x}" for x in swot_items(report, "opportunities")),
            "\n".join(f"• {x}" for x in swot_items(report, "threats")),
        ],
    ]
    t = Table(swot_data, colWidths=["50%", "50%"])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, 0), colors.HexColor("#d1fae5")),
        ("BACKGROUND", (1, 0), (1, 0), colors.HexColor("#fee2e2")),
        ("BACKGROUND", (0, 2), (0, 2), colors.HexColor("#dbeafe")),
        ("BACKGROUND", (1, 2), (1, 2), colors.HexColor("#fef3c7")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTNAME", (0, 2), (-1, 2), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("TEXTCOLOR", (0, 0), (-1, -1), dark),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#e5e7eb")),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("LEFTPADDING", (0, 0), (-1, -1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 4 * mm))

    # ── Financial Outlook ──────────────────────────────────────────────────
    story.append(Paragraph("Financial Outlook", H2))
    story.append(Paragraph(safe(report, "financial_outlook"), BODY))

    # ── ASEAN Opportunities ────────────────────────────────────────────────
    story.append(Paragraph("ASEAN Opportunities", H2))
    story.append(Paragraph(safe(report, "asean_opportunities"), BODY))

    # ── Risks ──────────────────────────────────────────────────────────────
    story.append(Paragraph("Key Risks", H2))
    story.append(Paragraph(safe(report, "risks"), BODY))

    # ── Recommendation ────────────────────────────────────────────────────
    story.append(Paragraph("Our Recommendation", H2))
    story.append(Paragraph(safe(report, "recommendation"), BODY))

    # ── Appendix: Data Sources ─────────────────────────────────────────────
    if refs:
        story.append(PageBreak())
        story.append(Paragraph("Appendix: Verified Data Sources", H1))
        story.append(HRFlowable(width="100%", color=indigo, thickness=0.5, spaceAfter=8))
        for ref in refs:
            story.append(Paragraph(
                f"<b>[{ref['id']}] {ref.get('source_name', 'Unknown')}</b>", BODY
            ))
            if ref.get("url"):
                story.append(Paragraph(
                    f"<link href='{ref['url']}'>{ref['url']}</link>", CITE
                ))
            story.append(Spacer(1, 2 * mm))

    doc.build(story)
    return buf.getvalue()


# ─── DOCX ────────────────────────────────────────────────────────────────────

def _build_docx(report: dict) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)

    def add_heading(text: str, level: int = 1):
        h = doc.add_heading(text, level=level)
        if level == 1:
            h.runs[0].font.color.rgb = RGBColor(*BRAND_COLOUR)
        return h

    def add_para(text: str):
        p = doc.add_paragraph(text)
        p.runs[0].font.size = Pt(11) if p.runs else None
        return p

    name = business_name(report)
    refs = references(report)

    # Cover
    add_heading(f"{BRAND} — Executive Blueprint", level=1)
    p = doc.add_paragraph(f"{name}  ·  Generated {datetime.now().strftime('%B %d, %Y')}")
    p.runs[0].font.color.rgb = RGBColor(107, 114, 128)
    p.runs[0].font.italic = True
    doc.add_paragraph(f"Confidence Score: {report.get('confidence_score', 0):.0f}%").runs[0].font.bold = True
    doc.add_paragraph()

    # Investor pitch
    if report.get("investor_pitch"):
        add_heading("Investor Pitch", level=2)
        p = doc.add_paragraph(f'"{safe(report, "investor_pitch")}"')
        p.runs[0].font.italic = True

    # Summary sections
    for key, title in [
        ("overall_summary", "Overview"),
        ("executive_summary", "Executive Summary"),
        ("financial_outlook", "Financial Outlook"),
        ("asean_opportunities", "ASEAN Opportunities"),
        ("risks", "Key Risks"),
        ("recommendation", "Our Recommendation"),
    ]:
        val = safe(report, key)
        if val:
            add_heading(title, level=2)
            add_para(val)

    # Key Findings
    kf = findings(report)
    if kf:
        add_heading("Key Findings", level=2)
        for f in kf:
            p = doc.add_paragraph(
                strip_citations(f.get("claim", "")),
                style="List Bullet",
            )
            lvl = f.get("confidence_level", "")
            run = p.add_run(f"  [{lvl}]")
            run.font.color.rgb = RGBColor(16, 185, 129) if lvl == "High" else RGBColor(245, 158, 11)
            run.font.bold = True

    # SWOT
    add_heading("SWOT Analysis", level=2)
    for quadrant, label in [
        ("strengths", "Strengths ✅"),
        ("weaknesses", "Weaknesses ⚠️"),
        ("opportunities", "Opportunities 🌍"),
        ("threats", "Threats 🔴"),
    ]:
        p = doc.add_paragraph()
        run = p.add_run(label)
        run.font.bold = True
        for item in swot_items(report, quadrant):
            doc.add_paragraph(item, style="List Bullet")

    # Appendix
    if refs:
        doc.add_page_break()
        add_heading("Appendix: Verified Data Sources", level=1)
        for ref in refs:
            p = doc.add_paragraph()
            run = p.add_run(f"[{ref['id']}] {ref.get('source_name', 'Unknown')}")
            run.font.bold = True
            if ref.get("url"):
                doc.add_paragraph(ref["url"]).runs[0].font.color.rgb = RGBColor(59, 130, 246)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── PPTX ────────────────────────────────────────────────────────────────────

def _build_pptx(report: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, RGBColor as PptxRGB

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank
    title_layout = prs.slide_layouts[0]

    indigo = PptxRGB(*BRAND_COLOUR)
    white = PptxRGB(255, 255, 255)
    dark = PptxRGB(30, 41, 59)
    amber = PptxRGB(245, 158, 11)
    green = PptxRGB(16, 185, 129)

    name = business_name(report)
    refs = references(report)
    kf = findings(report)

    def add_slide(title_text: str, body_text: str, notes_text: str = "",
                  bg_hex: str = "#0f172a"):
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        r, g, b = int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16)
        fill.fore_color.rgb = PptxRGB(r, g, b)

        # Brand watermark
        txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(3), Inches(0.3))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = BRAND.upper()
        run.font.size = Pt(8)
        run.font.color.rgb = indigo
        run.font.bold = True

        # Title
        txb_t = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(1.2))
        tf_t = txb_t.text_frame
        tf_t.word_wrap = True
        p = tf_t.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = white

        # Divider  
        line = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
            Inches(0.6), Inches(1.7), Inches(12), Emu(2),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = indigo
        line.line.fill.background()

        # Body
        if body_text:
            txb_b = slide.shapes.add_textbox(Inches(0.6), Inches(1.85), Inches(12), Inches(5.2))
            tf_b = txb_b.text_frame
            tf_b.word_wrap = True
            for i, line_text in enumerate(body_text.split("\n")):
                para = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
                run = para.add_run()
                run.text = line_text.strip()
                run.font.size = Pt(16 if line_text.startswith("•") or line_text.startswith("–") else 14)
                run.font.color.rgb = PptxRGB(209, 213, 219)

        # Notes (citations)
        if notes_text or refs:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes_text or ""
            if refs:
                notes_tf.text += "\n\nSources:\n" + "\n".join(
                    f"[{r['id']}] {r.get('source_name','')} — {r.get('url','')}"
                    for r in refs
                )

        return slide

    score = report.get("confidence_score", 0)

    # Slide 1 — Title
    pitch_short = textwrap.shorten(safe(report, "investor_pitch"), width=200, placeholder="…")
    add_slide(
        f"{name}",
        f"Executive Blueprint\n\n\"{pitch_short}\"\n\nConfidence Score: {score:.0f}%   ·   {datetime.now().strftime('%B %Y')}",
        bg_hex="#0f172a",
    )

    # Slide 2 — Problem & Solution (Executive Summary)
    summary_wrapped = textwrap.fill(safe(report, "executive_summary"), width=90)
    add_slide(
        "Business Summary",
        summary_wrapped,
        bg_hex="#0f172a",
    )

    # Slide 3 — SWOT
    def bullets(items: list[str], cap: int = 3) -> str:
        return "\n".join(f"• {textwrap.shorten(x, 80)}" for x in items[:cap])

    swot_body = (
        "STRENGTHS\n" + bullets(swot_items(report, "strengths")) +
        "\n\nWEAKNESSES\n" + bullets(swot_items(report, "weaknesses")) +
        "\n\nOPPORTUNITIES\n" + bullets(swot_items(report, "opportunities")) +
        "\n\nTHREATS\n" + bullets(swot_items(report, "threats"))
    )
    add_slide("SWOT Analysis", swot_body, bg_hex="#0f172a")

    # Slide 4 — Financials
    fin_lines = [
        textwrap.shorten(safe(report, "financial_outlook"), 300, placeholder="…"),
        "",
        "Key Findings:",
    ] + [f"• {textwrap.shorten(strip_citations(f.get('claim','')), 90)}" for f in kf[:4]]
    add_slide("Financial Outlook", "\n".join(fin_lines), bg_hex="#0f172a")

    # Slide 5 — AI Verdict
    rec = textwrap.fill(safe(report, "recommendation"), width=90)
    risk = textwrap.shorten(safe(report, "risks"), 220, placeholder="…")
    add_slide(
        "AI Verdict",
        f"Recommendation:\n{rec}\n\nKey Risks:\n{risk}",
        bg_hex="#0f172a",
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── XLSX ────────────────────────────────────────────────────────────────────

def _build_xlsx(report: dict, timeline: list[dict] | None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import (
        Alignment, Font, PatternFill, Border, Side
    )
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "Financial Timeline"

    indigo_hex = "6366F1"
    dark_hex = "0F172A"
    header_fill = PatternFill("solid", fgColor=indigo_hex)
    subheader_fill = PatternFill("solid", fgColor="1E293B")
    alt_fill = PatternFill("solid", fgColor="1A1F2E")
    red_fill = PatternFill("solid", fgColor="7F1D1D")

    header_font = Font(color="FFFFFF", bold=True, size=11, name="Calibri")
    white_font = Font(color="E2E8F0", size=10, name="Calibri")
    bold_white = Font(color="FFFFFF", bold=True, size=10, name="Calibri")
    dim_font = Font(color="94A3B8", size=9, name="Calibri")
    red_font = Font(color="FCA5A5", bold=True, size=10, name="Calibri")
    green_font = Font(color="6EE7B7", bold=True, size=10, name="Calibri")

    mid_align = Alignment(horizontal="center", vertical="center")
    left_align = Alignment(horizontal="left", vertical="center")

    # ── Sheet title ───────────────────────────────────────────────────────
    ws.merge_cells("A1:F1")
    cell = ws["A1"]
    cell.value = f"{BRAND} — Financial Sandbox: {business_name(report)}"
    cell.font = Font(color="FFFFFF", bold=True, size=14, name="Calibri")
    cell.fill = PatternFill("solid", fgColor=indigo_hex)
    cell.alignment = mid_align
    ws.row_dimensions[1].height = 32

    ws.merge_cells("A2:F2")
    cell = ws["A2"]
    cell.value = (
        f"Generated {datetime.now().strftime('%B %d, %Y')}  |  "
        f"Confidence: {report.get('confidence_score', 0):.0f}%  |  "
        "Import this file into QuickBooks, Excel, or Google Sheets"
    )
    cell.font = dim_font
    cell.fill = PatternFill("solid", fgColor="1E293B")
    cell.alignment = mid_align
    ws.row_dimensions[2].height = 18

    # ── Timeline header ───────────────────────────────────────────────────
    headers = ["Month", "Cash Balance (RM)", "Monthly Inflow (RM)", "Monthly Outflow (RM)", "Net (RM)", "Status"]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = mid_align
    ws.row_dimensions[4].height = 22

    rows = timeline or []
    for row_idx, pt in enumerate(rows, 5):
        cash = pt.get("cash", 0)
        inflow = pt.get("inflow", 0)
        outflow = pt.get("outflow", 0)
        net = inflow - outflow
        is_critical = cash < 0
        fill = red_fill if is_critical else (alt_fill if row_idx % 2 == 0 else subheader_fill)

        row_data = [
            f"Month {pt.get('month', row_idx - 4)}",
            cash,
            inflow,
            outflow,
            net,
            "⚠ Negative" if is_critical else "✓ Healthy",
        ]
        for col, val in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col, value=val)
            cell.fill = fill
            cell.alignment = mid_align
            if col == 1:
                cell.font = white_font
            elif col in (2, 3, 4, 5):
                cell.number_format = '#,##0'
                cell.font = red_font if is_critical else (green_font if col in (3, 5) and val >= 0 else white_font)
            else:
                cell.font = red_font if is_critical else Font(color="6EE7B7", size=10, name="Calibri")
            ws.row_dimensions[row_idx].height = 18

    if not rows:
        ws.merge_cells("A5:F6")
        cell = ws["A5"]
        cell.value = "No timeline data supplied. Run the Financial Runway Simulator and export again."
        cell.font = dim_font
        cell.alignment = mid_align

    # Column widths
    for i, w in enumerate([14, 22, 22, 24, 16, 14], 1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ── Financial Highlights summary ──────────────────────────────────────
    sum_row = len(rows) + 7
    ws.merge_cells(f"A{sum_row}:F{sum_row}")
    cell = ws.cell(row=sum_row, column=1, value="FINANCIAL HIGHLIGHTS")
    cell.font = header_font
    cell.fill = PatternFill("solid", fgColor=indigo_hex)
    cell.alignment = mid_align

    highlights = [
        ("Confidence Score", f"{report.get('confidence_score', 0):.0f}%"),
        ("Financial Outlook", textwrap.shorten(safe(report, "financial_outlook"), 120)),
        ("Recommendation", textwrap.shorten(safe(report, "recommendation"), 120)),
        ("Key Risks", textwrap.shorten(safe(report, "risks"), 120)),
    ]
    for offset, (label, value) in enumerate(highlights, sum_row + 1):
        ws.cell(row=offset, column=1, value=label).font = bold_white
        ws.cell(row=offset, column=1).fill = subheader_fill
        ws.cell(row=offset, column=1).alignment = left_align
        cell = ws.cell(row=offset, column=2, value=value)
        cell.font = white_font
        cell.fill = subheader_fill
        ws.merge_cells(f"B{offset}:F{offset}")
        ws.row_dimensions[offset].height = 16

    # ── Data Sources sheet ────────────────────────────────────────────────
    ws2 = wb.create_sheet("Data Sources")
    ws2.sheet_view.showGridLines = False

    ws2.merge_cells("A1:D1")
    cell = ws2["A1"]
    cell.value = "Appendix: Verified Data Sources"
    cell.font = Font(color="FFFFFF", bold=True, size=13, name="Calibri")
    cell.fill = PatternFill("solid", fgColor=indigo_hex)
    cell.alignment = mid_align
    ws2.row_dimensions[1].height = 28

    src_headers = ["#", "Source Name", "URL", "Confidence"]
    for col, h in enumerate(src_headers, 1):
        cell = ws2.cell(row=3, column=col, value=h)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = mid_align

    for ref_idx, ref in enumerate(references(report), 4):
        ws2.cell(row=ref_idx, column=1, value=f"[{ref.get('id','?')}]").font = bold_white
        ws2.cell(row=ref_idx, column=2, value=ref.get("source_name", "")).font = white_font
        cell = ws2.cell(row=ref_idx, column=3, value=ref.get("url", ""))
        cell.font = Font(color="60A5FA", size=10, name="Calibri")
        ws2.cell(row=ref_idx, column=4, value="Verified ✓").font = Font(color="6EE7B7", size=10, name="Calibri")
        for col in range(1, 5):
            ws2.cell(row=ref_idx, column=col).fill = (
                alt_fill if ref_idx % 2 == 0 else subheader_fill
            )
        ws2.row_dimensions[ref_idx].height = 16

    for col, w in enumerate([8, 30, 55, 14], 1):
        ws2.column_dimensions[get_column_letter(col)].width = w

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ─── Endpoint ────────────────────────────────────────────────────────────────

FORMAT_CONFIG = {
    "pdf":  ("application/pdf",                             "executive-blueprint.pdf"),
    "docx": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",
             "executive-blueprint.docx"),
    "pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation",
             "pitch-deck.pptx"),
    "xlsx": ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
             "financial-sandbox.xlsx"),
}


@router.post("")
async def export_blueprint(req: ExportRequest):
    fmt = req.format.lower()
    if fmt not in FORMAT_CONFIG:
        from fastapi import HTTPException
        raise HTTPException(400, detail=f"Unsupported format: {fmt}. Use pdf, docx, pptx, or xlsx.")

    if fmt == "pdf":
        data = _build_pdf(req.report)
    elif fmt == "docx":
        data = _build_docx(req.report)
    elif fmt == "pptx":
        data = _build_pptx(req.report)
    else:
        data = _build_xlsx(req.report, req.timeline)

    media_type, filename = FORMAT_CONFIG[fmt]
    return StreamingResponse(
        io.BytesIO(data),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
